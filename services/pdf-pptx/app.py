from flask import Flask, request, jsonify, send_file, send_from_directory, abort
from flask_cors import CORS
from werkzeug.exceptions import HTTPException
import io, os, urllib.parse, tempfile, shutil, errno, logging
from uuid import uuid4
from concurrent.futures import ThreadPoolExecutor
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu

app = Flask(__name__, static_folder="web", static_url_path="")
logging.basicConfig(level=logging.INFO)

# CORS
CORS(app, resources={r"/*": {"origins": [
    "http://localhost:5173",
    "https://77-tools.xyz",
    "https://www.77-tools.xyz",
    "https://popular-77.vercel.app"
], "expose_headers": ["Content-Disposition"], "methods": ["GET","POST","OPTIONS"], "allow_headers": ["Content-Type"]}})

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOADS_DIR = os.path.join(BASE_DIR, "uploads")
OUTPUTS_DIR = os.path.join(BASE_DIR, "outputs")
os.makedirs(UPLOADS_DIR, exist_ok=True)
os.makedirs(OUTPUTS_DIR, exist_ok=True)

executor = ThreadPoolExecutor(max_workers=2)
JOBS = {}

def safe_base_name(filename: str) -> str:
    base = os.path.splitext(os.path.basename(filename or "output"))[0]
    return base.replace("/", "").replace("\\", "").strip() or "output"

def set_progress(job_id, p, msg=None):
    info = JOBS.setdefault(job_id, {})
    info["progress"] = int(p)
    if msg is not None:
        info["message"] = msg

def send_download_memory(path: str, download_name: str, ctype: str):
    if not path or not os.path.exists(path):
        return jsonify({"error": "output file missing"}), 500
    with open(path, "rb") as f:
        data = f.read()
    resp = send_file(io.BytesIO(data), mimetype=ctype, as_attachment=True, download_name=download_name, conditional=False)
    resp.direct_passthrough = False
    resp.headers["Content-Length"] = str(len(data))
    resp.headers["Cache-Control"] = "no-store"
    resp.headers["Content-Disposition"] = f"attachment; filename*=UTF-8''{urllib.parse.quote(download_name)}"
    return resp

def safe_move(src: str, dst: str):
    os.makedirs(os.path.dirname(dst), exist_ok=True)
    try:
        if os.path.exists(dst): 
            os.remove(dst)
        os.replace(src, dst)
    except OSError as e:
        if getattr(e, "errno", None) == errno.EXDEV:
            shutil.copy2(src, dst)
            try: 
                os.unlink(src)
            except FileNotFoundError: 
                pass
        else:
            raise

def perform_pptx_conversion(in_path: str, base_name: str, scale: float = 1.0, job_id: str = None):
    """
    PDF → PPTX (페이지당 슬라이드 1장, 전체 이미지 맞춤)
    """
    with tempfile.TemporaryDirectory(dir=OUTPUTS_DIR) as tmp:
        doc = fitz.open(in_path)
        mat = fitz.Matrix(scale, scale)
        prs = Presentation()
        
        # 첫 페이지 이미지 크기에 맞춰 슬라이드 크기(EMU) 설정
        first = doc.load_page(0)
        pix0 = first.get_pixmap(matrix=mat, alpha=False)
        # 1 px ≈ 9525 EMU
        prs.slide_width = Emu(pix0.width * 9525)
        prs.slide_height = Emu(pix0.height * 9525)

        for i in range(doc.page_count):
            # job_id가 있을 때만 progress 업데이트 (비동기식에서만)
            if job_id:
                set_progress(job_id, 10 + int(80 * (i + 1) / doc.page_count), f"페이지 {i+1}/{doc.page_count} 처리 중")
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_path = os.path.join(tmp, f"{base_name}_{i+1:02d}.png")
            pix.save(img_path)

            blank = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(blank)
            # 전체 채우기
            slide.shapes.add_picture(img_path, Emu(0), Emu(0), width=prs.slide_width, height=prs.slide_height)

        final_name = f"{base_name}.pptx"
        final_path = os.path.join(OUTPUTS_DIR, final_name)
        if os.path.exists(final_path): 
            os.remove(final_path)
        prs.save(final_path)
        return final_path, final_name, "application/vnd.openxmlformats-officedocument.presentationml.presentation"

@app.get("/")
def index():
    return send_from_directory(app.static_folder, "index.html")

@app.route("/<path:path>")
def spa(path):
    # API는 제외
    if path in ("health", "convert", "convert-async") or path.startswith(("job/", "download/", "api/")):
        abort(404)
    full = os.path.join(app.static_folder, path)
    if os.path.isfile(full):
        return send_from_directory(app.static_folder, path)
    return send_from_directory(app.static_folder, "index.html")

@app.route("/health", methods=["GET"])
def health():
    return "ok", 200

@app.route("/convert-async", methods=["POST"])
def convert_async():
    f = request.files.get("file") or request.files.get("pdfFile")
    if not f: 
        return jsonify({"error":"file field is required"}), 400
    base_name = safe_base_name(f.filename)
    job_id = uuid4().hex
    in_path = os.path.join(UPLOADS_DIR, f"{job_id}.pdf")
    f.save(in_path)

    def clamp_num(v, lo, hi, default, T=float):
        try: 
            x = T(v)
        except: 
            return default
        return max(lo, min(hi, x))

    scale = clamp_num(request.form.get("scale","1.0"), 0.2, 2.0, 1.0, float)

    JOBS[job_id] = {"status":"pending","progress":1,"message":"대기 중"}
    app.logger.info(f"[{job_id}] uploaded: {in_path}, base={base_name}, scale={scale}")

    def run_job():
        try:
            set_progress(job_id, 10, "변환 준비 중")
            out_path, name, ctype = perform_pptx_conversion(in_path, base_name, scale=scale, job_id=job_id)
            JOBS[job_id] = {"status":"done","path":out_path,"name":name,"ctype":ctype,"progress":100,"message":"완료"}
        except Exception as e:
            app.logger.exception("convert error")
            JOBS[job_id] = {"status":"error","error":str(e),"progress":0,"message":"오류"}
        finally:
            try: 
                os.remove(in_path)
            except: 
                pass

    executor.submit(run_job)
    return jsonify({"job_id": job_id}), 202

@app.route("/job/<job_id>", methods=["GET"])
def job_status(job_id):
    info = JOBS.get(job_id)
    if not info: 
        return jsonify({"error":"job not found"}), 404
    info.setdefault("progress", 0)
    info.setdefault("message","")
    return jsonify(info), 200

@app.route("/download/<job_id>", methods=["GET"])
def job_download(job_id):
    info = JOBS.get(job_id)
    if not info: 
        return jsonify({"error":"job not found"}), 404
    if info.get("status") != "done": 
        return jsonify({"error":"not ready"}), 409
    return send_download_memory(info["path"], info["name"], info["ctype"])

# API aliases: /api/pdf-pptx/*
@app.route("/api/pdf-pptx/health", methods=["GET"])
def _a1(): 
    return health()

@app.route("/api/pdf-pptx/convert-async", methods=["POST"])
def _a2(): 
    return convert_async()

@app.route("/api/pdf-pptx/job/<job_id>", methods=["GET"])
def _a3(job_id): 
    return job_status(job_id)

@app.route("/api/pdf-pptx/download/<job_id>", methods=["GET"])
def _a4(job_id): 
    return job_download(job_id)

@app.post("/convert")
def convert_sync():
    f = request.files.get("file") or request.files.get("pdfFile")
    if not f:
        return jsonify({"error":"file field is required"}), 400

    base_name = safe_base_name(f.filename)
    job_id = uuid4().hex
    in_path = os.path.join(UPLOADS_DIR, f"{job_id}.pdf")
    f.save(in_path)

    # doc과 동일: quality는 받되, PPTX 변환에는 사용하지 않음(무시)
    quality = request.form.get("quality", "low")  # "low"(빠른) | "standard"
    # scale은 UI에서 숨기므로 기본 1.0(옵션으로 전달되면 반영)
    def clamp_num(v, lo, hi, default, T=float):
        try: 
            x = T(v)
        except: 
            return default
        return max(lo, min(hi, x))
    scale = clamp_num(request.form.get("scale","1.0"), 0.2, 2.0, 1.0, float)

    try:
        out_path, name, ctype = perform_pptx_conversion(in_path, base_name, scale=scale, job_id=None)
        return send_download_memory(out_path, name, ctype)
    except Exception as e:
        app.logger.exception("convert sync error")
        return jsonify({"error": str(e)}), 500
    finally:
        try: 
            os.remove(in_path)
        except: 
            pass

# /api 별칭(온렌더에서 /api 경로도 동작하도록)
@app.post("/api/pdf-pptx/convert")
def _alias_convert_sync():
    return convert_sync()

@app.errorhandler(Exception)
def handle_exception(e):
    if isinstance(e, HTTPException):
        return e
    app.logger.exception("Unhandled exception")
    return jsonify({"error": "Internal server error"}), 500

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)