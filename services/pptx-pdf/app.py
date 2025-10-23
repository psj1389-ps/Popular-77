from flask import Flask, request, jsonify, send_file, send_from_directory, abort, render_template, make_response
from flask_cors import CORS
from werkzeug.exceptions import HTTPException
import io, os, urllib.parse, tempfile, shutil, errno, logging, subprocess, shlex
from uuid import uuid4
from concurrent.futures import ThreadPoolExecutor
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu
from urllib.parse import quote
import unicodedata
from threading import Semaphore

# LibreOffice 기반 변환 사용
ALLOWED_EXTS = {".pptx", ".ppt", ".odp"}

app = Flask(__name__, template_folder="templates", static_folder="static")
logging.basicConfig(level=logging.INFO)

# 최대 업로드 크기 설정 (60MB)
app.config["MAX_CONTENT_LENGTH"] = int(os.getenv("MAX_CONTENT_LENGTH_MB", "60")) * 1024 * 1024

# 동시성 제어 (메모리 사용량 제한)
SEMA = Semaphore(int(os.getenv("MAX_CONCURRENCY", "1")))

# 타임아웃 설정
app.config["SEND_FILE_MAX_AGE_DEFAULT"] = 0

# 요청 로깅 추가
@app.before_request
def _trace():
    app.logger.info(f">>> {request.method} {request.path}")

# CORS
CORS(app, resources={r"/*": {"origins": [
    "http://localhost:5173",
    "https://77-tools.xyz",
    "https://www.77-tools.xyz",
    "https://popular-77.vercel.app",
    "https://popular-77-xbqq.onrender.com"
], "expose_headers": ["Content-Disposition"], "methods": ["GET","POST","OPTIONS"], "allow_headers": ["Content-Type"], "supports_credentials": False}})

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOADS_DIR = os.path.join(BASE_DIR, "uploads")
OUTPUTS_DIR = os.path.join(BASE_DIR, "outputs")
os.makedirs(UPLOADS_DIR, exist_ok=True)
os.makedirs(OUTPUTS_DIR, exist_ok=True)

executor = ThreadPoolExecutor(max_workers=2)
JOBS = {}

def ascii_fallback(name: str) -> str:
    """유니코드 파일명을 ASCII로 변환하여 안전한 파일명 생성"""
    a = unicodedata.normalize("NFKD", name).encode("ascii", "ignore").decode("ascii") or "converted.pdf"
    return "".join(c for c in a if c.isalnum() or c in ".- ") or "converted.pdf"

def _set_pdf_disposition(resp, pdf_name: str):
    """RFC 5987 + ASCII fallback으로 Content-Disposition 헤더 설정"""
    ascii_name = ascii_fallback(pdf_name)
    resp.headers["Content-Disposition"] = (
        f'attachment; filename="{ascii_name}"; filename*=UTF-8\'\'{quote(pdf_name)}'
    )
    resp.headers["Cache-Control"] = "no-store"
    return resp

def safe_base_name(filename: str) -> str:
    base = os.path.splitext(os.path.basename(filename or "output"))[0]
    return base.replace("/", "").replace("\\", "").strip() or "output"

def set_progress(job_id, p, msg=None):
    info = JOBS.setdefault(job_id, {})
    info["progress"] = int(p)
    if msg is not None:
        info["message"] = msg

def send_download_memory(path: str, download_name: str, ctype: str):
    """메모리 효율적인 파일 다운로드 - 스트리밍 사용"""
    try:
        # 파일 존재 확인
        if not os.path.exists(path):
            app.logger.error(f"전송할 파일이 존재하지 않음: {path}")
            return jsonify({"error": "파일을 찾을 수 없습니다"}), 404
        
        file_size = os.path.getsize(path)
        if file_size == 0:
            app.logger.error(f"전송할 파일이 비어있음: {path}")
            return jsonify({"error": "파일이 비어있습니다"}), 500
        
        app.logger.info(f"파일 전송 시작: {download_name} ({file_size} bytes)")
        
        # 스트리밍 응답으로 메모리 사용량 최적화
        resp = send_file(
            path, 
            mimetype=ctype, 
            as_attachment=True, 
            download_name=download_name,
            conditional=True  # 조건부 요청 지원
        )
        
        # RFC5987 호환 헤더 설정
        ascii_name = ascii_fallback(download_name)
        resp.headers["Content-Disposition"] = (
            f'attachment; filename="{ascii_name}"; filename*=UTF-8\'\'{urllib.parse.quote(download_name)}'
        )
        resp.headers["Content-Length"] = str(file_size)
        resp.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
        resp.headers["Pragma"] = "no-cache"
        resp.headers["Expires"] = "0"
        
        app.logger.info(f"파일 전송 완료: {download_name}")
        return resp
    except Exception as e:
        app.logger.error(f"파일 전송 실패: {path} - {e}")
        return jsonify({"error": f"파일 전송 실패: {str(e)}"}), 500

def safe_move(src: str, dst: str):
    os.makedirs(os.path.dirname(dst), exist_ok=True)
    try:
        if os.path.exists(dst): 
            os.remove(dst)
        os.replace(src, dst)
    except OSError as e:
        if getattr(e, "errno", None) == errno.EXDEV:
            shutil.copy2(src, dst)
            try: 
                os.unlink(src)
            except FileNotFoundError: 
                pass
        else:
            raise

def perform_libreoffice(in_path: str, out_pdf_path: str):
    import shlex, subprocess, shutil, os
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice: 
        raise RuntimeError("LibreOffice not found")
    
    outdir = os.path.dirname(out_pdf_path)
    os.makedirs(outdir, exist_ok=True)
    
    ext = os.path.splitext(in_path)[1].lower()
    filters = {
        ".ppt": "impress_pdf_Export", ".pptx": "impress_pdf_Export",
        ".xls": "calc_pdf_Export", ".xlsx": "calc_pdf_Export", 
        ".doc": "writer_pdf_Export", ".docx": "writer_pdf_Export"
    }
    conv = f"pdf:{filters.get(ext,'pdf')}"
    
    cmd = f'{shlex.quote(soffice)} --headless --nologo --nofirststartwizard --convert-to {conv} --outdir {shlex.quote(outdir)} {shlex.quote(in_path)}'
    app.logger.info(f"[LO] cmd={cmd}")
    
    p = subprocess.run(cmd, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    app.logger.info(f"[LO] rc={p.returncode}, out={p.stdout.decode('utf-8','ignore')[:200]}")
    
    if p.returncode != 0:
        raise RuntimeError(p.stderr.decode("utf-8","ignore"))
    
    produced = os.path.join(outdir, os.path.splitext(os.path.basename(in_path))[0] + ".pdf")
    if not os.path.exists(produced):
        raise RuntimeError(f"PDF not produced: {produced}")
    
    # 최종 경로로 이동
    if produced != out_pdf_path:
        if os.path.exists(out_pdf_path):
            os.remove(out_pdf_path)
        os.replace(produced, out_pdf_path)

# 이미지 기반 변환 (기존 PDF to PPTX 기능 유지)

def perform_pptx_conversion(in_path: str, base_name: str, scale: float = 1.0, job_id: str = None):
    """
    PDF → PPTX (페이지당 슬라이드 1장, 전체 이미지 맞춤)
    """
    with tempfile.TemporaryDirectory(dir=OUTPUTS_DIR) as tmp:
        doc = None
        try:
            # 파일 존재 여부 확인
            if not os.path.exists(in_path):
                raise FileNotFoundError(f"PDF 파일을 찾을 수 없습니다: {in_path}")
            
            # 파일 크기 확인
            file_size = os.path.getsize(in_path)
            if file_size == 0:
                raise ValueError(f"PDF 파일이 비어있습니다: {in_path}")
            
            app.logger.info(f"PDF 파일 로딩 시작: {in_path} ({file_size} bytes)")
            
            # PDF 파일 열기 시도
            try:
                doc = fitz.open(in_path)
            except Exception as pdf_error:
                app.logger.error(f"PDF 파일 열기 실패: {str(pdf_error)}")
                raise ValueError(f"PDF 문서를 로드하지 못했습니다. 파일이 손상되었거나 유효하지 않은 PDF 형식일 수 있습니다: {str(pdf_error)}")
            
            # PDF 문서 유효성 검증
            if doc is None:
                raise ValueError("PDF 문서 객체가 None입니다")
            
            if doc.is_closed:
                raise ValueError("PDF 문서가 이미 닫혀있습니다")
            
            page_count = doc.page_count
            if page_count == 0:
                raise ValueError("PDF 문서에 페이지가 없습니다")
            
            app.logger.info(f"PDF 문서 로딩 성공: {page_count}페이지")
            
            mat = fitz.Matrix(scale, scale)
            prs = Presentation()
            
            # 첫 페이지 이미지 크기에 맞춰 슬라이드 크기(EMU) 설정
            try:
                first = doc.load_page(0)
                pix0 = first.get_pixmap(matrix=mat, alpha=False)
                # 1 px ≈ 9525 EMU
                prs.slide_width = Emu(pix0.width * 9525)
                prs.slide_height = Emu(pix0.height * 9525)
                app.logger.info(f"슬라이드 크기 설정: {pix0.width}x{pix0.height} px")
            except Exception as page_error:
                app.logger.error(f"첫 페이지 로딩 실패: {str(page_error)}")
                raise ValueError(f"PDF 첫 페이지를 처리할 수 없습니다: {str(page_error)}")

            for i in range(page_count):
                try:
                    # job_id가 있을 때만 progress 업데이트 (비동기식에서만)
                    if job_id:
                        set_progress(job_id, 10 + int(80 * (i + 1) / page_count), f"페이지 {i+1}/{page_count} 처리 중")
                    
                    page = doc.load_page(i)
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    img_path = os.path.join(tmp, f"{base_name}_{i+1:02d}.png")
                    pix.save(img_path)
                    
                    # 이미지 파일 생성 확인
                    if not os.path.exists(img_path):
                        raise ValueError(f"페이지 {i+1} 이미지 생성 실패: {img_path}")

                    blank = prs.slide_layouts[6]  # Blank
                    slide = prs.slides.add_slide(blank)
                    # 전체 채우기
                    slide.shapes.add_picture(img_path, Emu(0), Emu(0), width=prs.slide_width, height=prs.slide_height)
                    
                    app.logger.debug(f"페이지 {i+1}/{page_count} 처리 완료")
                    
                except Exception as page_error:
                    app.logger.error(f"페이지 {i+1} 처리 실패: {str(page_error)}")
                    raise ValueError(f"페이지 {i+1} 처리 중 오류가 발생했습니다: {str(page_error)}")

            final_name = f"{base_name}.pdf"
            final_path = os.path.join(OUTPUTS_DIR, final_name)
            if os.path.exists(final_path): 
                os.remove(final_path)
            
            # PDF 파일 저장
            try:
                prs.save(final_path)
                app.logger.info(f"PDF 파일 저장 완료: {final_path}")
            except Exception as save_error:
                app.logger.error(f"PDF 파일 저장 실패: {str(save_error)}")
                raise ValueError(f"PDF 파일 저장 중 오류가 발생했습니다: {str(save_error)}")
            
            # 파일이 제대로 생성되었는지 확인
            if not os.path.exists(final_path):
                raise Exception(f"PDF 파일 생성 실패: {final_path}")
            
            file_size = os.path.getsize(final_path)
            if file_size == 0:
                raise Exception(f"PDF 파일이 비어있음: {final_path}")
            
            app.logger.info(f"PDF 변환 완료: {final_name} ({file_size} bytes)")
            return final_path, final_name, "application/pdf"
            
        except FileNotFoundError as e:
            app.logger.error(f"파일 없음 오류: {str(e)}")
            raise Exception(f"파일을 찾을 수 없습니다: {str(e)}")
        except ValueError as e:
            app.logger.error(f"PDF 처리 오류: {str(e)}")
            raise Exception(str(e))
        except Exception as e:
            app.logger.error(f"PPTX 변환 오류: {str(e)}")
            raise Exception(f"PPTX 변환 실패: {str(e)}")
        finally:
            if doc is not None and not doc.is_closed:
                try:
                    doc.close()
                    app.logger.debug("PDF 문서 닫기 완료")
                except Exception as close_error:
                    app.logger.warning(f"PDF 문서 닫기 실패: {str(close_error)}")

@app.get("/")
def index():
    # 이 서비스 전용 텍스트/수치(템플릿에서 사용)
    page = {
        "title": "PDF to PPTX Converter",
        "subtitle": "PDF 문서를 PowerPoint로 안정적으로 변환",
        "accept": ".pdf",
        "max_mb": os.getenv("MAX_CONTENT_LENGTH_MB", "100"),
        "service": "pptx-pdf"
    }
    resp = make_response(render_template("index.html", page=page))
    resp.headers["Cache-Control"] = "no-store"
    resp.headers["Pragma"] = "no-cache"
    return resp

@app.route("/<path:path>")
def spa(path):
    # API는 제외
    if path in ("health", "convert", "convert-async") or path.startswith(("job/", "download/", "api/")):
        abort(404)
    full = os.path.join(app.static_folder, path)
    if os.path.isfile(full):
        return send_from_directory(app.static_folder, path)
    return send_from_directory(app.static_folder, "index.html")

@app.route("/health", methods=["GET"])
def health():
    import sys, shutil, subprocess
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    lo_ver = "not_found"
    if soffice:
        v = subprocess.run([soffice, "--version"], capture_output=True, text=True)
        lo_ver = (v.stdout or v.stderr or "").strip()
    return {
        "service": "pptx-pdf",
        "engine": "libreoffice",
        "soffice_path": soffice,
        "libreoffice": lo_ver,
        "pdfservices_sdk": None,  # Adobe 완전 제거 확인
        "allowed_exts": sorted(list(ALLOWED_EXTS)),
        "python": sys.version,
    }

@app.route("/convert-async", methods=["POST"])
def convert_async():
    f = request.files.get("file") or request.files.get("pdfFile")
    if not f: 
        return jsonify({"error":"file field is required"}), 400
    base_name = safe_base_name(f.filename)
    job_id = uuid4().hex
    in_path = os.path.join(UPLOADS_DIR, f"{job_id}.pdf")
    
    try:
        f.save(in_path)
        
        # 업로드된 파일 확인
        if not os.path.exists(in_path):
            return jsonify({"error": "파일 업로드 실패"}), 500
        
        upload_size = os.path.getsize(in_path)
        if upload_size == 0:
            return jsonify({"error": "업로드된 파일이 비어있음"}), 400
        
        app.logger.info(f"비동기 파일 업로드 완료: {f.filename} ({upload_size} bytes)")
        
    except Exception as e:
        return jsonify({"error": f"파일 업로드 오류: {str(e)}"}), 500

    def clamp_num(v, lo, hi, default, T=float):
        try: 
            x = T(v)
        except: 
            return default
        return max(lo, min(hi, x))

    scale = clamp_num(request.form.get("scale","1.0"), 0.2, 2.0, 1.0, float)

    JOBS[job_id] = {"status":"pending","progress":1,"message":"대기 중"}
    app.logger.info(f"[{job_id}] uploaded: {in_path}, base={base_name}, scale={scale}")

    def run_job():
        try:
            set_progress(job_id, 30, "이미지 기반 변환 중")
            out_path, name, ctype = perform_pptx_conversion(in_path, base_name, scale=scale, job_id=job_id)
            set_progress(job_id, 90, "파일 저장 중")
        except Exception as e:
            app.logger.error(f"변환 실패: {str(e)}")
            JOBS[job_id] = {"status":"error","error":f"변환 실패: {str(e)}","progress":0,"message":"변환 실패"}
            return
        
        JOBS[job_id] = {"status":"done","path":out_path,"name":name,"ctype":ctype,"progress":100,"message":"완료"}
        app.logger.info(f"비동기 변환 완료: {job_id} -> {name}")
        
        try: 
            if os.path.exists(in_path):
                os.remove(in_path)
        except Exception as e:
            app.logger.warning(f"비동기 임시 파일 삭제 실패: {str(e)}")

    executor.submit(run_job)
    return jsonify({"job_id": job_id}), 202

@app.route("/job/<job_id>", methods=["GET"])
def job_status(job_id):
    info = JOBS.get(job_id)
    if not info: 
        return jsonify({"error":"job not found"}), 404
    info.setdefault("progress", 0)
    info.setdefault("message","")
    
    # 에러 상태인 경우 상세 정보 제공
    if info.get("status") == "error":
        error_msg = info.get("error", "알 수 없는 오류")
        app.logger.error(f"Job {job_id} 상태 확인 - 오류: {error_msg}")
        return jsonify({
            "status": "error",
            "error": error_msg,
            "progress": info.get("progress", 0),
            "message": info.get("message", "")
        }), 200
    
    return jsonify(info), 200

@app.route("/download/<job_id>", methods=["GET"])
def job_download(job_id):
    info = JOBS.get(job_id)
    if not info: 
        return jsonify({"error":"job not found"}), 404
    if info.get("status") != "done": 
        return jsonify({"error":"not ready"}), 409
    return send_download_memory(info["path"], info["name"], info["ctype"])

# API aliases: /api/pdf-pptx/*
@app.post("/upload")
@app.post("/convert")
def convert_sync():
    # 세마포어로 동시성 제어
    if not SEMA.acquire(timeout=float(os.getenv("QUEUE_TIMEOUT", "0"))):
        return jsonify({"error": "busy"}), 503
    
    try:
        f = request.files.get("file") or request.files.get("document")
        if not f:
            return jsonify({"error": "file field is required"}), 400

        in_name = f.filename or "input.pptx"  # 서비스별 기본값
        base, ext = os.path.splitext(os.path.basename(in_name))
        ext = ext.lower()
        
        # 디버깅 로그 추가
        app.logger.info(f"Original filename: {f.filename}")
        app.logger.info(f"Processed in_name: {in_name}")
        app.logger.info(f"Extracted base: {base}")
        app.logger.info(f"Extracted ext: {ext}")
        
        if ext not in ALLOWED_EXTS:
            return jsonify({"error": f"unsupported extension: {ext}"}), 415

        jid = uuid4().hex
        in_path = os.path.join(UPLOADS_DIR, f"{jid}{ext}")
        out_name = f"{base}.pdf"
        tmp_out = os.path.join(OUTPUTS_DIR, f"{jid}.pdf")
        final_out = os.path.join(OUTPUTS_DIR, out_name)
        
        # 최종 출력 파일명 로그
        app.logger.info(f"Final output name: {out_name}")

        f.save(in_path)
        try:
            perform_libreoffice(in_path, tmp_out)
            if not _is_pdf(tmp_out):
                raise RuntimeError("Output is not a valid PDF")
            if tmp_out != final_out:
                if os.path.exists(final_out): os.remove(final_out)
                os.replace(tmp_out, final_out)

            size = os.path.getsize(final_out)
            # 스트리밍으로 직접 전송(메모리 절약)
            resp = send_file(
                final_out,
                as_attachment=True,
                download_name=out_name,         # 기본 파일명(추가로 아래에서 RFC5987 설정)
                mimetype="application/pdf",
                conditional=False
            )
            resp.headers["Content-Length"] = str(size)
            return _set_pdf_disposition(resp, out_name)
        except Exception as e:
            app.logger.exception("conversion failed")
            return jsonify({"error": str(e)}), 500
        finally:
            try: os.remove(in_path)
            except: pass
    finally:
        SEMA.release()

# /api aliases (frontend uses /api paths only)
@app.route("/api/pdf-pptx/health", methods=["GET", "HEAD"])
def _pptx_a_health(): return health()

@app.route("/api/pdf-pptx/convert-async", methods=["POST","OPTIONS"])
def _pptx_a_convert_async(): return convert_async()

@app.route("/api/pdf-pptx/convert", methods=["POST","OPTIONS"])
def _pptx_a_convert_sync(): return convert_sync()

@app.route("/api/pdf-pptx/job/<job_id>", methods=["GET","HEAD"])
def _pptx_a_job(job_id): return job_status(job_id)

@app.route("/api/pdf-pptx/download/<job_id>", methods=["GET","HEAD"])
def _pptx_a_download(job_id): return job_download(job_id)

@app.errorhandler(Exception)
def handle_exception(e):
    if isinstance(e, HTTPException):
        return e
    
    # 상세한 오류 로깅
    import traceback
    app.logger.error(f"Unhandled exception: {str(e)}")
    app.logger.error(f"Exception type: {type(e).__name__}")
    app.logger.error(f"Traceback: {traceback.format_exc()}")
    
    # 개발 모드에서는 상세 오류 정보 반환
    if app.debug:
        return jsonify({
            "error": "Internal server error",
            "details": str(e),
            "type": type(e).__name__
        }), 500
    else:
        return jsonify({"error": "Internal server error"}), 500

def _is_pdf(path: str) -> bool:
    try:
        with open(path, "rb") as f: 
            return f.read(4) == b"%PDF"
    except: 
        return False

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port, debug=True)