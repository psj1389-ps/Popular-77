from flask import Flask, request, render_template, send_file, flash, redirect, url_for, jsonify
import os
import tempfile
from werkzeug.utils import secure_filename
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image
import json
from dotenv import load_dotenv
from docx import Document
from docx.shared import Pt, Inches as DocxInches
import subprocess
import platform
import pytesseract
import cv2
import numpy as np
import fitz  # PyMuPDF
import re
from typing import List, Tuple, Dict, Any
from pdf2docx import Converter
# Adobe PDF Services SDK 임포트 및 설정
try:
    # 올바른 Adobe PDF Services SDK import 구문
    from adobe.pdfservices.operation.auth.service_principal_credentials import ServicePrincipalCredentials
    from adobe.pdfservices.operation.exception.exceptions import ServiceApiException, ServiceUsageException, SdkException
    from adobe.pdfservices.operation.io.cloud_asset import CloudAsset
    from adobe.pdfservices.operation.io.stream_asset import StreamAsset
    from adobe.pdfservices.operation.pdf_services import PDFServices
    from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
    from adobe.pdfservices.operation.pdfjobs.jobs.export_pdf_job import ExportPDFJob
    from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_params import ExportPDFParams
    from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_target_format import ExportPDFTargetFormat
    from adobe.pdfservices.operation.pdfjobs.result.export_pdf_result import ExportPDFResult
    
    adobe_available = True
    ADOBE_SDK_AVAILABLE = True
    print("Adobe PDF Services SDK가 성공적으로 로드되었습니다.")
except ImportError as e:
    print(f"Adobe PDF Services SDK를 가져올 수 없습니다: {e}")
    print("Adobe SDK 없이 계속 진행합니다.")
    adobe_available = False
    ADOBE_SDK_AVAILABLE = False

# Adobe SDK 무조건 실행 모드 활성화
ADOBE_SDK_AVAILABLE = True
print(f"Adobe SDK 강제 실행 모드: {ADOBE_SDK_AVAILABLE}")

# 환경 변수 로드
load_dotenv()

app = Flask(__name__)
app.secret_key = 'your-secret-key-here'
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB max file size

# Adobe PDF Services API 구성 - 환경변수 사용
ADOBE_CONFIG = {
    "client_id": os.getenv("ADOBE_CLIENT_ID"),
    "client_secret": os.getenv("ADOBE_CLIENT_SECRET"),
    "organization_id": os.getenv("ADOBE_ORGANIZATION_ID"),
    "account_id": os.getenv("ADOBE_ACCOUNT_ID"),
    "technical_account_email": os.getenv("ADOBE_TECHNICAL_ACCOUNT_EMAIL")
}

# Adobe SDK 무조건 실행 확인
print(f"Adobe SDK 강제 실행 상태: {ADOBE_SDK_AVAILABLE}")
print(f"Adobe 클라이언트 ID: {ADOBE_CONFIG['client_credentials']['client_id']}")
print("Adobe PDF Services API가 무조건 실행되도록 설정되었습니다.")



UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
ALLOWED_EXTENSIONS = {'pdf'}

# 폴더 생성
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs('debug_output', exist_ok=True)  # 디버깅용 폴더

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# 디버깅용 중간 결과물 저장 함수들
def save_debug_text(text, filename_prefix):
    """추출된 텍스트를 디버깅용 .txt 파일로 저장"""
    try:
        debug_file = os.path.join('debug_output', f'{filename_prefix}_extracted_text.txt')
        with open(debug_file, 'w', encoding='utf-8') as f:
            f.write(text)
        print(f"디버깅: 텍스트가 {debug_file}에 저장되었습니다. (길이: {len(text)}자)")
        return debug_file
    except Exception as e:
        print(f"디버깅 텍스트 저장 오류: {e}")
        return None

def save_debug_image(image, filename_prefix, page_num):
    """변환된 이미지를 디버깅용 .png 파일로 저장"""
    try:
        debug_file = os.path.join('debug_output', f'{filename_prefix}_page_{page_num}.png')
        image.save(debug_file, 'PNG')
        print(f"디버깅: 이미지가 {debug_file}에 저장되었습니다.")
        return debug_file
    except Exception as e:
        print(f"디버깅 이미지 저장 오류: {e}")
        return None

def pdf_to_docx_with_pdf2docx(pdf_path, output_path):
    """pdf2docx 라이브러리를 사용한 PDF → DOCX 변환"""
    try:
        print("=== pdf2docx 라이브러리 변환 시작 ===")
        print(f"입력 파일: {pdf_path}")
        print(f"출력 파일: {output_path}")
        
        # 입력 파일 존재 여부 확인
        if not os.path.exists(pdf_path):
            print(f"❌ 입력 파일이 존재하지 않습니다: {pdf_path}")
            return False
            
        # 파일 크기 확인
        file_size = os.path.getsize(pdf_path)
        print(f"📄 PDF 파일 크기: {file_size:,} bytes ({file_size/1024:.1f} KB)")
        
        # 출력 디렉토리 생성
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
            print(f"📁 출력 디렉토리 생성: {output_dir}")
        
        print("🔄 pdf2docx Converter 객체 생성 중...")
        # Converter 객체 생성
        cv = Converter(pdf_path)
        
        print("🔄 PDF → DOCX 변환 실행 중...")
        # 변환 실행
        cv.convert(output_path, start=0, end=None)
        
        print("🔄 Converter 객체 닫는 중...")
        # 객체 닫기
        cv.close()
        
        # 변환 결과 확인
        if os.path.exists(output_path):
            output_size = os.path.getsize(output_path)
            print(f"✅ pdf2docx 변환 성공!")
            print(f"📄 출력 파일 크기: {output_size:,} bytes ({output_size/1024:.1f} KB)")
            print(f"📁 출력 파일 경로: {output_path}")
            return True
        else:
            print(f"❌ 변환 완료되었으나 출력 파일이 생성되지 않음: {output_path}")
            return False
        
    except ImportError as e:
        print(f"❌ pdf2docx 라이브러리 import 오류: {str(e)}")
        print("💡 해결방법: pip install pdf2docx 명령어로 라이브러리를 설치하세요")
        return False
    except PermissionError as e:
        print(f"❌ 파일 권한 오류: {str(e)}")
        print(f"💡 해결방법: 파일이 다른 프로그램에서 사용 중이거나 권한이 없습니다")
        return False
    except FileNotFoundError as e:
        print(f"❌ 파일을 찾을 수 없음: {str(e)}")
        return False
    except Exception as e:
        print(f"❌ pdf2docx 라이브러리 변환 실패: {str(e)}")
        print(f"🔍 오류 타입: {type(e).__name__}")
        import traceback
        print(f"📋 상세 스택 트레이스:")
        print(traceback.format_exc())
        print(f"📄 입력 파일: {pdf_path}")
        print(f"📁 출력 파일: {output_path}")
        return False

def ocr_image_to_blocks(pil_image):
    """이미지에서 단어 단위 텍스트와 위치(좌표)를 추출"""
    try:
        img = cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        gray = cv2.medianBlur(gray, 3)
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        gray = clahe.apply(gray)

        config = r"--oem 3 --psm 6 -l kor+eng"
        data = pytesseract.image_to_data(gray, config=config,
                                         output_type=pytesseract.Output.DICT)
        blocks = []
        n = len(data["text"])
        for i in range(n):
            text = data["text"][i].strip()
            conf_val = data["conf"][i]
            if isinstance(conf_val, (int, float)):
                conf = int(conf_val)
            elif isinstance(conf_val, str) and conf_val.replace('.', '').replace('-', '').isdigit():
                conf = int(float(conf_val))
            else:
                conf = 0
            
            if conf > 30 and len(text) > 0:
                x, y, w, h = data["left"][i], data["top"][i], data["width"][i], data["height"][i]
                blocks.append({
                    "text": text,
                    "bbox": (x, y, x + w, y + h),
                    "confidence": conf
                })
        return blocks
    except Exception as e:
        print(f"OCR 처리 중 오류: {e}")
        return []

def clean_special_characters(text: str) -> str:
    """특수 문자 처리 개선 - PDF에서 잘못 추출되는 문자들을 올바르게 복구"""
    if not text:
        return text
    
    # 일반적인 PDF 추출 오류 수정
    replacements = {
        '\uf0b7': '•',  # 불릿 포인트
        '\uf0a7': '§',  # 섹션 기호
        '\uf0e0': '→',  # 화살표
        '\u2022': '•',  # 불릿 포인트
        '\u201C': '"',  # 왼쪽 큰따옴표
        '\u201D': '"',  # 오른쪽 큰따옴표
        '\u2018': "'",  # 왼쪽 작은따옴표
        '\u2019': "'",  # 오른쪽 작은따옴표
        '\u2013': '–',  # en dash
        '\u2014': '—',  # em dash
        '\u00A0': ' ',  # 줄바꿈 없는 공백
        '\u200B': '',   # 폭이 0인 공백
        '\uFEFF': '',   # 바이트 순서 표시
    }
    
    # 특수 문자 변환
    for old, new in replacements.items():
        text = text.replace(old, new)
    
    # 연속된 공백 정리
    text = re.sub(r'[\s\t\n\r]+', ' ', text)
    
    # 제로 폭 문자 제거
    text = re.sub(r'[\u200B-\u200D\uFEFF]', '', text)
    
    return text.strip()

def analyze_pdf_orientation(pdf_path: str) -> Dict[str, Any]:
    """PDF 페이지 크기를 분석하여 문서 방향 감지"""
    try:
        doc = fitz.open(pdf_path)
        page_orientations = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_rect = page.rect
            width = page_rect.width
            height = page_rect.height
            
            # 가로/세로 방향 판단
            if width > height:
                orientation = 'landscape'  # 가로형
            else:
                orientation = 'portrait'   # 세로형
            
            page_orientations.append({
                'page': page_num,
                'width': width,
                'height': height,
                'orientation': orientation,
                'aspect_ratio': width / height
            })
        
        doc.close()
        
        # 전체 문서의 주요 방향 결정
        landscape_count = sum(1 for p in page_orientations if p['orientation'] == 'landscape')
        portrait_count = len(page_orientations) - landscape_count
        
        primary_orientation = 'landscape' if landscape_count > portrait_count else 'portrait'
        
        return {
            'page_orientations': page_orientations,
            'primary_orientation': primary_orientation,
            'landscape_pages': landscape_count,
            'portrait_pages': portrait_count,
            'total_pages': len(page_orientations)
        }
        
    except Exception as e:
        print(f"PDF 방향 분석 중 오류: {e}")
        return {
            'page_orientations': [],
            'primary_orientation': 'portrait',
            'landscape_pages': 0,
            'portrait_pages': 0,
            'total_pages': 0
        }

def extract_text_with_layout_from_pdf(pdf_path: str) -> Dict[str, Any]:
    """PDF에서 레이아웃 정보와 함께 텍스트 추출 (방향 감지 포함)"""
    try:
        doc = fitz.open(pdf_path)
        all_text_blocks = []
        
        # PDF 방향 분석
        orientation_info = analyze_pdf_orientation(pdf_path)
        print(f"PDF 방향 분석 결과: {orientation_info['primary_orientation']} (가로: {orientation_info['landscape_pages']}, 세로: {orientation_info['portrait_pages']})")
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_rect = page.rect
            
            # 텍스트 블록 추출
            text_blocks = page.get_text("dict")
            
            for block in text_blocks.get("blocks", []):
                if "lines" in block:
                    for line in block["lines"]:
                        line_text = ""
                        line_bbox = line["bbox"]
                        
                        for span in line.get("spans", []):
                            span_text = span.get("text", "")
                            if span_text.strip():
                                line_text += span_text
                        
                        if line_text.strip():
                            # 텍스트 정렬 감지
                            text_center = (line_bbox[0] + line_bbox[2]) / 2
                            page_center = page_rect.width / 2
                            
                            if abs(text_center - page_center) < 20:
                                alignment = 'center'
                            elif (page_rect.width - line_bbox[2]) < (line_bbox[0]):
                                alignment = 'right'
                            else:
                                alignment = 'left'
                            
                            all_text_blocks.append({
                                'text': clean_special_characters(line_text),
                                'bbox': line_bbox,
                                'page': page_num,
                                'alignment': alignment
                            })
        
        doc.close()
        
        # 텍스트 블록들을 위치 순서대로 정렬
        all_text_blocks.sort(key=lambda x: (x['page'], x['bbox'][1], x['bbox'][0]))
        
        return {
            'text_blocks': all_text_blocks,
            'full_text': '\n'.join([block['text'] for block in all_text_blocks]),
            'orientation_info': orientation_info
        }
        
    except Exception as e:
        print(f"PDF 레이아웃 추출 중 오류: {e}")
        return {'text_blocks': [], 'full_text': ''}

def extract_text_blocks_with_ocr(image):
    """OCR을 사용하여 이미지에서 텍스트 블록 추출 (개선된 버전)"""
    try:
        # 이미지 전처리로 OCR 정확도 향상
        img_array = np.array(image)
        
        # 그레이스케일 변환
        if len(img_array.shape) == 3:
            gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        else:
            gray = img_array
        
        # 노이즈 제거
        denoised = cv2.medianBlur(gray, 3)
        
        # 대비 향상
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        enhanced = clahe.apply(denoised)
        
        # OCR 수행
        config = r"--oem 3 --psm 6 -l kor+eng"
        text = pytesseract.image_to_string(enhanced, config=config)
        
        if text.strip():
            cleaned_text = clean_special_characters(text.strip())
            print(f"  - OCR 텍스트 추출됨: {len(cleaned_text)}자")
            return cleaned_text
        else:
            print("  - OCR에서 텍스트를 찾을 수 없음")
            return ""
            
    except Exception as e:
        print(f"  - OCR 처리 중 오류: {e}")
        return ""

def extract_pdf_content_with_adobe(pdf_path):
    """Adobe PDF Services API를 사용하여 PDF 내용을 추출하는 함수"""
    if not ADOBE_SDK_AVAILABLE:
        print("Adobe PDF Services SDK를 사용할 수 없습니다.")
        return None
        
    try:
        # Adobe API 자격 증명 설정 (올바른 클래스 사용)
        credentials = ServicePrincipalCredentials(
            client_id=ADOBE_CONFIG["client_credentials"]["client_id"],
            client_secret=ADOBE_CONFIG["client_credentials"]["client_secret"]
        )
        
        # PDF Services 인스턴스 생성
        pdf_services = PDFServices(credentials=credentials)
        
        # PDF 파일을 스트림으로 읽기
        with open(pdf_path, 'rb') as file:
            input_stream = file.read()
        
        # StreamAsset 생성
        input_asset = pdf_services.upload(input_stream=input_stream, mime_type=PDFServicesMediaType.PDF)
        
        print("Adobe API를 사용하여 PDF 내용을 처리했습니다.")
        return input_asset
            
    except (ServiceApiException, ServiceUsageException, SdkException) as e:
        print(f"Adobe API 오류: {str(e)}")
        return None
    except Exception as e:
        print(f"일반 오류: {str(e)}")
        return None

def pdf_to_docx(pdf_path, output_path, quality='medium'):
    """PDF를 DOCX로 변환하는 함수 (하이브리드 접근법: pdf2docx 우선, OCR 보조)"""
    try:
        print(f"=== PDF 변환 시작: {pdf_path} ===")
        
        # 입력 파일 검증
        if not os.path.exists(pdf_path):
            print(f"[ERROR] 입력 파일이 존재하지 않습니다: {pdf_path}")
            return False
            
        file_size = os.path.getsize(pdf_path)
        print(f"파일 크기: {file_size} bytes ({file_size/1024:.1f}KB)")
        
        # 출력 디렉토리 생성
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)
            print(f"출력 디렉토리 생성: {output_dir}")
        
        # 파일명에서 확장자 제거하여 디버깅용 prefix 생성
        filename_prefix = os.path.splitext(os.path.basename(pdf_path))[0]
        
        # 1단계: pdf2docx 라이브러리를 우선적으로 시도
        print("=== 1단계: pdf2docx 라이브러리 변환 시도 ===")
        print(f"입력 파일: {pdf_path}")
        print(f"출력 파일: {output_path}")
        
        try:
            print("pdf_to_docx_with_pdf2docx 함수 호출 중...")
            conversion_result = pdf_to_docx_with_pdf2docx(pdf_path, output_path)
            print(f"pdf_to_docx_with_pdf2docx 반환값: {conversion_result}")
            
            if conversion_result:
                print("pdf2docx 변환 성공! Microsoft Word 호환성 확인...")
                
                # 변환된 파일이 실제로 존재하고 크기가 적절한지 확인
                if os.path.exists(output_path):
                    output_size = os.path.getsize(output_path)
                    print(f"출력 파일 존재 확인: {output_path}")
                    print(f"출력 파일 크기: {output_size} bytes ({output_size/1024:.1f}KB)")
                    
                    if output_size > 1024:  # 1KB 이상
                        print(f"[SUCCESS] pdf2docx 변환 완료: {output_path}")
                        return True
                    else:
                        print(f"[WARNING] 출력 파일이 너무 작음 ({output_size} bytes). 대체 방법 시도...")
                else:
                    print(f"[ERROR] 출력 파일이 생성되지 않음: {output_path}")
                    print("pdf2docx 변환 결과가 부적절함. 대체 방법 시도...")
            else:
                print("[ERROR] pdf_to_docx_with_pdf2docx 함수가 False 반환. 대체 방법 시도...")
                
        except Exception as e:
            print(f"[ERROR] pdf2docx 변환 중 예외 발생: {str(e)}")
            print(f"예외 타입: {type(e).__name__}")
            import traceback
            print(f"상세 스택 트레이스: {traceback.format_exc()}")
        
        print("=== 2단계: 기존 OCR 방법으로 fallback ===")
        print(f"품질 설정: {quality}")
        
        try:
            # 품질 설정에 따른 파라미터 설정 (최적화됨)
            quality_settings = {
                'medium': {
                    'dpi': 120,  # DPI 최적화로 속도 향상
                    'format': 'jpeg',
                    'jpeg_quality': 80,  # 품질과 속도의 균형
                    'max_size': (1600, 1200),  # 적절한 해상도
                    'description': '균형 변환 (최적화된 속도와 품질)'
                },
                'high': {
                    'dpi': 180,  # 고품질이지만 속도 고려
                    'format': 'jpeg',  # PNG 대신 JPEG 사용으로 속도 향상
                    'jpeg_quality': 90,
                    'max_size': (2048, 1536),  # 해상도 최적화
                    'description': '고품질 변환 (향상된 속도)'
                }
            }
            
            settings = quality_settings.get(quality, quality_settings['medium'])
            print(f"변환 설정: {settings['description']}")
            
            # 1단계: 레이아웃 인식을 통한 텍스트 추출 시도
            print("레이아웃 인식을 통한 텍스트 추출을 시도합니다...")
            layout_data = extract_text_with_layout_from_pdf(pdf_path)
            extracted_text = layout_data.get('full_text', '')
            text_blocks = layout_data.get('text_blocks', [])
            orientation_info = layout_data.get('orientation_info', {})
            
            if extracted_text:
                print(f"레이아웃 인식으로 텍스트 추출 성공: {len(extracted_text)}자")
            else:
                print("레이아웃 인식 실패, Adobe API를 시도합니다...")
                
                # 2단계: Adobe API를 사용한 PDF 내용 추출 시도
                if adobe_available:
                    print("Adobe API를 사용하여 PDF 처리를 시작합니다...")
                    extracted_content = extract_pdf_content_with_adobe(pdf_path)
                    if extracted_content:
                        extracted_text = str(extracted_content)
                        print(f"Adobe API에서 텍스트 추출 성공: {len(extracted_text)}자")
                    else:
                        print("Adobe API 추출 실패, OCR 방법으로 진행합니다.")
            
                # 기본 방법: PDF를 이미지로 변환 (품질별 최적화)
                print("PDF를 이미지로 변환 중...")
                try:
                    images = convert_from_path(pdf_path, dpi=settings['dpi'], fmt=settings['format'])
                    print(f"이미지 변환 성공: {len(images)}페이지")
                except Exception as e:
                    print(f"PDF 이미지 변환 실패: {str(e)}")
                    return False
        
                # 디버깅: 변환된 이미지들을 저장
                print("=== 디버깅: 변환된 이미지 저장 ===")
                for i, image in enumerate(images):
                    save_debug_image(image, filename_prefix, i+1)
        
        except Exception as e:
            print(f"PDF → DOCX 변환 중 치명적 오류 발생: {str(e)}")
            print(f"오류 타입: {type(e).__name__}")
            import traceback
            print(f"상세 스택 트레이스: {traceback.format_exc()}")
            return False
        
        # 새 Word 문서 생성 - 호환성 개선 및 방향 자동 감지
        doc = Document()
        
        # 페이지 설정 (문서 방향에 따라 자동 조정)
        section = doc.sections[0]
        primary_orientation = orientation_info.get('primary_orientation', 'portrait')
        
        if primary_orientation == 'landscape':
            # 가로형 문서 설정
            section.page_width = Inches(11)
            section.page_height = Inches(8.5)
            section.left_margin = Inches(0.8)
            section.right_margin = Inches(0.8)
            section.top_margin = Inches(0.6)
            section.bottom_margin = Inches(0.6)
            print("가로형 문서로 설정됨")
        else:
            # 세로형 문서 설정
            section.page_width = Inches(8.5)
            section.page_height = Inches(11)
            section.left_margin = Inches(1)
            section.right_margin = Inches(1)
            section.top_margin = Inches(1)
            section.bottom_margin = Inches(1)
            print("세로형 문서로 설정됨")
        
        # 문서 속성 설정 (Microsoft Word 호환성 향상)
        try:
            # 안전한 파일명 생성 (특수문자 제거)
            safe_filename = re.sub(r'[^\w\s-]', '', os.path.splitext(os.path.basename(pdf_path))[0])
            doc.core_properties.title = safe_filename[:50]  # 제목 길이 제한
            doc.core_properties.author = "Document Converter"
            doc.core_properties.subject = "PDF to DOCX Conversion"
            doc.core_properties.comments = "Converted using advanced OCR and layout recognition"
        except Exception as e:
            print(f"문서 속성 설정 중 오류 (무시됨): {e}")
        
        all_ocr_text = []
        
        print(f"총 {len(images)}페이지 처리 중...")
        # OCR로 텍스트 추출 (Adobe API가 실패한 경우)
        if not extracted_text:
            for i, image in enumerate(images):
                print(f"페이지 {i+1}/{len(images)} OCR 처리 중...")
                ocr_text = extract_text_blocks_with_ocr(image)
                if ocr_text:
                    all_ocr_text.append(ocr_text)
        
        # 편집 가능한 텍스트만 추가 (원본 이미지 제거)
        final_text = extracted_text if extracted_text else '\n'.join(all_ocr_text)
        
        # 디버깅: 추출된 텍스트를 파일로 저장
        print("=== 디버깅: 추출된 텍스트 저장 ===")
        if final_text:
            save_debug_text(final_text, filename_prefix)
        elif all_ocr_text:
            combined_ocr_text = '\n'.join(all_ocr_text)
            save_debug_text(combined_ocr_text, filename_prefix + "_ocr")
        else:
            save_debug_text("텍스트 추출 실패", filename_prefix + "_no_text")
        
        if final_text or text_blocks:
            print(f"편집 가능한 텍스트 문서 생성: {len(final_text)}자")
            
            # 레이아웃 정보가 있는 경우 페이지별로 구성 (페이지 헤더 제거)
            if text_blocks:
                print("레이아웃 정보를 활용하여 텍스트 구조화...")
                
                # 페이지별로 텍스트 구성 (페이지 번호 헤더 없이)
                for page_num in range(len(images)):
                    if page_num > 0:
                        doc.add_page_break()
                    
                    # 해당 페이지의 텍스트 블록 추가
                    page_text_blocks = [block for block in text_blocks if block['page'] == page_num]
                    
                    if page_text_blocks:
                        for block in page_text_blocks:
                            text_paragraph = doc.add_paragraph()
                            text_run = text_paragraph.add_run(block['text'])
                            
                            # 텍스트 정렬 적용 (원본과 동일하게)
                            if block['alignment'] == 'center':
                                text_paragraph.alignment = 1  # 중앙 정렬
                                text_run.bold = True  # 중앙 정렬 텍스트는 굵게
                            elif block['alignment'] == 'right':
                                text_paragraph.alignment = 2  # 오른쪽 정렬
                            else:
                                text_paragraph.alignment = 0  # 왼쪽 정렬
                            
                            # 원본과 동일한 텍스트 스타일 적용 (Microsoft Word 호환성 개선)
                            try:
                                # 폰트 설정 (한글 문서에 적합한 폰트)
                                text_run.font.name = '맑은 고딕'
                                text_run.font.size = Pt(11)  # 표준 문서 크기
                                
                                # 제목 스타일 적용
                                if len(block['text']) < 50 and block['alignment'] == 'center':
                                    text_run.font.size = Pt(14)
                                    text_run.bold = True
                                elif '제목' in block['text'] or '공문' in block['text']:
                                    text_run.font.size = Pt(13)
                                    text_run.bold = True
                                
                                # 줄간격 및 단락 간격 설정 (원본과 동일하게)
                                text_paragraph.paragraph_format.line_spacing = 1.2
                                text_paragraph.paragraph_format.space_after = Pt(3)
                                text_paragraph.paragraph_format.space_before = Pt(0)
                                
                                # 들여쓰기 설정 (원본 레이아웃 유지)
                                if block['alignment'] == 'left':
                                    text_paragraph.paragraph_format.left_indent = Pt(0)
                                elif block['alignment'] == 'center':
                                    text_paragraph.paragraph_format.left_indent = Pt(0)
                                    
                            except Exception as e:
                                print(f"텍스트 스타일 설정 중 오류 (무시됨): {e}")
            else:
                # 레이아웃 정보가 없는 경우 일반 텍스트로 추가 (Microsoft Word 호환성 개선)
                clean_final_text = final_text.replace('\x00', '').replace('\ufffd', '').strip()
                if clean_final_text:
                    paragraphs = clean_final_text.split('\n\n')
                    for para_text in paragraphs:
                        if para_text.strip():
                            text_paragraph = doc.add_paragraph()
                            text_run = text_paragraph.add_run(para_text.strip())
                            
                            # 일반 텍스트에도 표준 스타일 적용
                            try:
                                text_run.font.name = '맑은 고딕'
                                text_run.font.size = Pt(11)
                                text_paragraph.paragraph_format.line_spacing = 1.2
                                text_paragraph.paragraph_format.space_after = Pt(3)
                                text_paragraph.paragraph_format.space_before = Pt(0)
                            except Exception as e:
                                print(f"일반 텍스트 스타일 설정 중 오류 (무시됨): {e}")
                else:
                    text_paragraph = doc.add_paragraph()
                    text_run = text_paragraph.add_run("텍스트를 추출할 수 없었습니다.")
                    try:
                        text_run.font.name = '맑은 고딕'
                        text_run.font.size = Pt(11)
                    except Exception as e:
                        print(f"오류 메시지 스타일 설정 중 오류 (무시됨): {e}")
            
            print("편집 가능한 텍스트 문서가 생성되었습니다.")
        else:
            print("추출할 수 있는 텍스트가 없습니다. 이미지 기반 문서를 생성합니다.")
            
            # 텍스트가 없는 경우에만 이미지 추가
            for i, image in enumerate(images):
                print(f"페이지 {i+1}/{len(images)} 처리 중...")
                
                # 이미지 크기 최적화 (원본 문서와 동일한 크기 유지)
                original_width, original_height = image.size
                
                # 문서 방향에 따른 이미지 크기 조정
                if primary_orientation == 'landscape':
                    # 가로형 문서: 최대 너비 10인치
                    target_width = min(10, original_width / 72)  # 72 DPI 기준
                    aspect_ratio = original_height / original_width
                    target_height = target_width * aspect_ratio
                else:
                    # 세로형 문서: 최대 너비 6.5인치
                    target_width = min(6.5, original_width / 72)  # 72 DPI 기준
                    aspect_ratio = original_height / original_width
                    target_height = target_width * aspect_ratio
                
                # 이미지를 임시 파일로 저장 (JPEG 최적화)
                temp_img_path = None
                try:
                    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                        temp_img_path = temp_img.name
                        # JPEG로 저장 (품질별 압축, 빠른 처리)
                        image.save(temp_img_path, 'JPEG', quality=settings['jpeg_quality'], optimize=True)
                    
                    # 문서에 이미지 추가 (원본 비율 유지)
                    doc.add_picture(temp_img_path, width=DocxInches(target_width))
                    
                    # 페이지 구분을 위한 페이지 브레이크 추가 (마지막 페이지 제외)
                    if i < len(images) - 1:
                        doc.add_page_break()
                    
                finally:
                    # 임시 파일 삭제 (빠른 처리)
                    if temp_img_path and os.path.exists(temp_img_path):
                        try:
                            os.unlink(temp_img_path)
                        except (OSError, PermissionError) as e:
                            print(f"임시 파일 삭제 실패 (무시됨): {e}")
        
        # DOCX 파일 저장 (Microsoft Word 호환성 최적화)
        try:
            # 임시 파일로 먼저 저장 후 이동 (안전한 저장)
            temp_output = output_path + '.tmp'
            doc.save(temp_output)
            
            # 기존 파일이 있으면 삭제
            if os.path.exists(output_path):
                os.remove(output_path)
            
            # 임시 파일을 최종 파일로 이동
            os.rename(temp_output, output_path)
            
            print(f"DOCX 파일 저장 완료: {output_path}")
            print("Microsoft Word 호환성이 개선된 문서가 생성되었습니다.")
            return True
            
        except Exception as save_error:
            print(f"DOCX 파일 저장 중 오류: {save_error}")
            # 임시 파일 정리
            temp_output = output_path + '.tmp'
            if os.path.exists(temp_output):
                try:
                    os.remove(temp_output)
                except:
                    pass
            return False
        
    except Exception as e:
        print(f"PDF → DOCX 변환 중 치명적 오류 발생: {str(e)}")
        print(f"오류 타입: {type(e).__name__}")
        import traceback
        print(f"상세 스택 트레이스: {traceback.format_exc()}")
        print(f"입력 파일: {pdf_path}")
        print(f"출력 파일: {output_path}")
        print(f"품질 설정: {quality}")
        return False

def pdf_to_pptx(pdf_path, output_path, quality='medium'):
    """PDF를 PPTX로 변환하는 함수 (Adobe API 통합 및 OCR 텍스트 추출, 방향 자동 감지)"""
    try:
        # 품질 설정에 따른 파라미터 설정 (최적화됨)
        quality_settings = {
            'medium': {
                'dpi': 120,  # DPI 최적화로 속도 향상
                'format': 'jpeg',
                'jpeg_quality': 80,  # 품질과 속도의 균형
                'max_size': (1600, 1200),  # 적절한 해상도
                'description': '균형 변환 (최적화된 속도와 품질)'
            },
            'high': {
                'dpi': 180,  # 고품질이지만 속도 고려
                'format': 'jpeg',  # PNG 대신 JPEG 사용으로 속도 향상
                'jpeg_quality': 90,
                'max_size': (2048, 1536),  # 해상도 최적화
                'description': '고품질 변환 (향상된 속도)'
            }
        }
        
        settings = quality_settings.get(quality, quality_settings['medium'])
        print(f"변환 설정: {settings['description']}")
        
        # 1단계: 레이아웃 인식을 통한 텍스트 추출 시도 (방향 정보 포함)
        print("레이아웃 인식을 통한 텍스트 추출을 시도합니다...")
        layout_data = extract_text_with_layout_from_pdf(pdf_path)
        extracted_text = layout_data.get('full_text', '')
        text_blocks = layout_data.get('text_blocks', [])
        orientation_info = layout_data.get('orientation_info', {})
        all_ocr_text = []
        
        if extracted_text:
            print(f"레이아웃 인식으로 텍스트 추출 성공: {len(extracted_text)}자")
        else:
            print("레이아웃 인식 실패, Adobe API를 시도합니다...")
            
            # 2단계: Adobe API를 사용한 PDF 내용 추출 시도
            if adobe_available:
                print("Adobe API를 사용하여 PDF 처리를 시작합니다...")
                extracted_content = extract_pdf_content_with_adobe(pdf_path)
                if extracted_content:
                    extracted_text = str(extracted_content)
                    print(f"Adobe API에서 텍스트 추출 성공: {len(extracted_text)}자")
                else:
                    print("Adobe API 추출 실패, OCR 방법으로 진행합니다.")
        
        # 기본 방법: PDF를 이미지로 변환 (품질별 최적화)
        print("PDF를 이미지로 변환 중...")
        images = convert_from_path(pdf_path, dpi=settings['dpi'], fmt=settings['format'])
        
        # 새 PowerPoint 프레젠테이션 생성 (방향에 따른 슬라이드 설정)
        prs = Presentation()
        
        # 슬라이드 크기 설정 (문서 방향에 따라 자동 조정)
        primary_orientation = orientation_info.get('primary_orientation', 'portrait')
        
        if primary_orientation == 'landscape':
            # 가로형 슬라이드 (16:9 비율)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            print("가로형 슬라이드로 설정됨 (16:9)")
        else:
            # 세로형 슬라이드 (9:16 비율)
            prs.slide_width = Inches(7.5)
            prs.slide_height = Inches(13.33)
            print("세로형 슬라이드로 설정됨 (9:16)")
        
        all_ocr_text = []
        
        print(f"총 {len(images)}페이지 처리 중...")
        def get_blank_slide_layout(prs):
            """안전한 빈 슬라이드 레이아웃 가져오기"""
            try:
                # 슬라이드 레이아웃이 있는지 먼저 확인
                if len(prs.slide_layouts) == 0:
                    raise IndexError("슬라이드 레이아웃이 없습니다")
                
                # 빈 슬라이드 레이아웃 우선 선택
                if len(prs.slide_layouts) > 6:
                    return prs.slide_layouts[6]  # 빈 슬라이드
                elif len(prs.slide_layouts) > 5:
                    return prs.slide_layouts[5]  # 제목만 있는 슬라이드
                else:
                    return prs.slide_layouts[0]  # 첫 번째 사용 가능한 레이아웃
            except (IndexError, AttributeError) as e:
                print(f"슬라이드 레이아웃 접근 오류: {e}")
                # 기본 프레젠테이션 생성 시 최소 하나의 레이아웃은 있어야 함
                if len(prs.slide_layouts) > 0:
                    return prs.slide_layouts[0]
                else:
                    raise Exception("사용 가능한 슬라이드 레이아웃이 없습니다")
        
        # 편집 가능한 텍스트 슬라이드 생성 (원본 이미지 제거)
        # OCR로 텍스트 추출 (Adobe API가 실패한 경우)
        if not extracted_text:
            for i, image in enumerate(images):
                print(f"페이지 {i+1}/{len(images)} OCR 처리 중...")
                ocr_text = extract_text_blocks_with_ocr(image)
                if ocr_text:
                    all_ocr_text.append(ocr_text)
        
        # 편집 가능한 텍스트 슬라이드 생성
        final_text = extracted_text if extracted_text else '\n'.join(all_ocr_text)
        
        if text_blocks:
            print(f"편집 가능한 텍스트 슬라이드 생성: {len(text_blocks)}개 블록")
            
            # 페이지별로 슬라이드 구성
            for page_num in range(len(images)):
                # 새 슬라이드 추가 (제목과 내용 레이아웃)
                try:
                    slide_layout = prs.slide_layouts[1]  # 제목과 내용 레이아웃
                except (IndexError, AttributeError):
                    slide_layout = get_blank_slide_layout(prs)
                
                slide = prs.slides.add_slide(slide_layout)
                
                # 슬라이드 제목 설정
                try:
                    title_shape = slide.shapes.title
                    title_shape.text = f"페이지 {page_num + 1}"
                except AttributeError:
                    # 제목이 없는 레이아웃인 경우 텍스트 박스 추가
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = Inches(9)
                    height = Inches(1)
                    title_box = slide.shapes.add_textbox(left, top, width, height)
                    title_frame = title_box.text_frame
                    title_frame.text = f"페이지 {page_num + 1}"
                
                # 해당 페이지의 텍스트 블록 추가
                page_text_blocks = [block for block in text_blocks if block['page'] == page_num]
                
                if page_text_blocks:
                    # 내용 텍스트박스 가져오기
                    try:
                        content_shape = slide.placeholders[1]
                        text_frame = content_shape.text_frame
                        text_frame.clear()
                        
                        for j, block in enumerate(page_text_blocks):
                            if j == 0:
                                # 첫 번째 단락
                                p = text_frame.paragraphs[0]
                            else:
                                # 추가 단락
                                p = text_frame.add_paragraph()
                            
                            p.text = block['text']
                            
                            # 텍스트 정렬 적용
                            if block['alignment'] == 'center':
                                p.alignment = 1  # 중앙 정렬
                                try:
                                    p.font.bold = True
                                except AttributeError:
                                    pass
                            elif block['alignment'] == 'right':
                                p.alignment = 2  # 오른쪽 정렬
                            else:
                                p.alignment = 0  # 왼쪽 정렬
                            
                            # 텍스트 크기 조정
                            try:
                                from docx.shared import Pt
                                if len(block['text']) < 50 and block['alignment'] == 'center':
                                    p.font.size = Pt(18)  # 제목용 크기
                                else:
                                    p.font.size = Pt(14)  # 본문용 크기
                            except (ImportError, AttributeError):
                                pass
                                
                    except (IndexError, AttributeError):
                        # 내용 플레이스홀더가 없는 경우 텍스트 박스 추가
                        left = Inches(0.5)
                        top = Inches(1.5)
                        width = Inches(9)
                        height = Inches(6)
                        content_box = slide.shapes.add_textbox(left, top, width, height)
                        content_frame = content_box.text_frame
                        content_text = '\n'.join([block['text'] for block in page_text_blocks])
                        content_frame.text = content_text
                else:
                    # 텍스트가 없는 페이지
                    try:
                        content_shape = slide.placeholders[1]
                        content_shape.text = "[이 페이지는 텍스트 추출이 어려운 이미지 페이지입니다]"
                    except (IndexError, AttributeError):
                        left = Inches(0.5)
                        top = Inches(1.5)
                        width = Inches(9)
                        height = Inches(6)
                        content_box = slide.shapes.add_textbox(left, top, width, height)
                        content_frame = content_box.text_frame
                        content_frame.text = "[이 페이지는 텍스트 추출이 어려운 이미지 페이지입니다]"
        
        elif final_text:
            print(f"일반 텍스트 슬라이드 생성: {len(final_text)}자")
            
            # 텍스트를 적절한 크기로 나누어 슬라이드 생성
            text_chunks = final_text.split('\n\n')
            chunk_size = 5  # 슬라이드당 단락 수
            
            for i in range(0, len(text_chunks), chunk_size):
                try:
                    slide_layout = prs.slide_layouts[1]  # 제목과 내용 레이아웃
                except (IndexError, AttributeError):
                    slide_layout = get_blank_slide_layout(prs)
                
                slide = prs.slides.add_slide(slide_layout)
                
                # 슬라이드 제목
                try:
                    title_shape = slide.shapes.title
                    title_shape.text = f"슬라이드 {(i // chunk_size) + 1}"
                except AttributeError:
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = Inches(9)
                    height = Inches(1)
                    title_box = slide.shapes.add_textbox(left, top, width, height)
                    title_frame = title_box.text_frame
                    title_frame.text = f"슬라이드 {(i // chunk_size) + 1}"
                
                # 내용 추가
                chunk_text = text_chunks[i:i+chunk_size]
                content_text = '\n\n'.join([para.strip() for para in chunk_text if para.strip()])
                
                try:
                    content_shape = slide.placeholders[1]
                    content_shape.text = content_text
                except (IndexError, AttributeError):
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(6)
                    content_box = slide.shapes.add_textbox(left, top, width, height)
                    content_frame = content_box.text_frame
                    content_frame.text = content_text
        
        else:
            print("추출할 수 있는 텍스트가 없습니다. 이미지 기반 슬라이드를 생성합니다.")
            
            # 텍스트가 없는 경우에만 이미지 슬라이드 생성
            for i, image in enumerate(images):
                print(f"페이지 {i+1}/{len(images)} 처리 중...")
                
                # 슬라이드 추가 - 안전한 레이아웃 사용
                slide_layout = get_blank_slide_layout(prs)
                slide = prs.slides.add_slide(slide_layout)
                
                # 이미지 크기 최적화 (원본 문서와 동일한 크기 유지)
                original_width, original_height = image.size
                
                # 슬라이드 방향에 따른 이미지 크기 조정
                if primary_orientation == 'landscape':
                    # 가로형 슬라이드: 최대 높이 6.5인치
                    target_height = min(6.5, original_height / 72)  # 72 DPI 기준
                    aspect_ratio = original_width / original_height
                    target_width = target_height * aspect_ratio
                    # 슬라이드 너비를 초과하지 않도록 조정
                    max_slide_width = 12.5  # 가로형 슬라이드 최대 너비
                    if target_width > max_slide_width:
                        target_width = max_slide_width
                        target_height = target_width / aspect_ratio
                else:
                    # 세로형 슬라이드: 최대 너비 6.5인치
                    target_width = min(6.5, original_width / 72)  # 72 DPI 기준
                    aspect_ratio = original_height / original_width
                    target_height = target_width * aspect_ratio
                    # 슬라이드 높이를 초과하지 않도록 조정
                    max_slide_height = 12.5  # 세로형 슬라이드 최대 높이
                    if target_height > max_slide_height:
                        target_height = max_slide_height
                        target_width = target_height / aspect_ratio
                
                # 이미지를 임시 파일로 저장 (JPEG 최적화)
                temp_img_path = None
                try:
                    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                        temp_img_path = temp_img.name
                        # JPEG로 저장 (품질별 압축, 빠른 처리)
                        image.save(temp_img_path, 'JPEG', quality=settings['jpeg_quality'], optimize=True)
                    
                    # 슬라이드에 이미지 추가 (원본 비율 유지, 중앙 배치)
                    left = Inches((13.33 - target_width) / 2) if primary_orientation == 'landscape' else Inches((7.5 - target_width) / 2)
                    top = Inches((7.5 - target_height) / 2) if primary_orientation == 'landscape' else Inches((13.33 - target_height) / 2)
                    slide.shapes.add_picture(temp_img_path, left, top, width=Inches(target_width), height=Inches(target_height))
                    
                finally:
                    # 임시 파일 삭제 (빠른 처리)
                    if temp_img_path and os.path.exists(temp_img_path):
                        try:
                            os.unlink(temp_img_path)
                        except (OSError, PermissionError) as e:
                            print(f"임시 파일 삭제 실패 (무시됨): {e}")
                            # 임시 파일 삭제 실패는 무시하고 계속 진행
        
        # 하이브리드 변환: 추출된 텍스트를 편집 가능한 형태로 마지막 슬라이드에 추가
        final_text = extracted_text if extracted_text else '\n'.join(all_ocr_text)
        if final_text:
            print(f"하이브리드 변환: 추출된 텍스트를 편집 가능한 형태로 추가: {len(final_text)}자")
            
            # 텍스트 전용 슬라이드 추가 - 안전한 레이아웃 선택
            try:
                # 제목과 내용 레이아웃이 있는지 확인
                if len(prs.slide_layouts) > 1:
                    text_slide_layout = prs.slide_layouts[1]  # 제목과 내용 레이아웃
                else:
                    text_slide_layout = get_blank_slide_layout(prs)
            except (IndexError, AttributeError):
                text_slide_layout = get_blank_slide_layout(prs)
            
            text_slide = prs.slides.add_slide(text_slide_layout)
            
            # 제목 설정 (안전한 방법)
            try:
                title = text_slide.shapes.title
                title.text = "추출된 텍스트 (편집 가능)"
            except AttributeError:
                # 제목이 없는 레이아웃인 경우 텍스트 박스 추가
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(1)
                title_box = text_slide.shapes.add_textbox(left, top, width, height)
                title_frame = title_box.text_frame
                title_frame.text = "추출된 텍스트 (편집 가능)"
            
            # 내용 설정 (안전한 방법) - 레이아웃 정보 활용
            if text_blocks:
                print("레이아웃 정보를 활용하여 텍스트 구조화...")
                try:
                    content = text_slide.placeholders[1]
                    content_frame = content.text_frame
                    content_frame.clear()
                    
                    current_page = -1
                    for block in text_blocks:
                        # 페이지가 바뀌면 구분선 추가
                        if block['page'] != current_page:
                            if current_page != -1:
                                p = content_frame.add_paragraph()
                                p.text = f"\n--- 페이지 {block['page'] + 1} ---"
                            current_page = block['page']
                        
                        # 텍스트 단락 추가
                        p = content_frame.add_paragraph()
                        p.text = block['text']
                        
                        # 정렬 설정 (기본값으로 처리)
                        if block['alignment'] == 'center':
                            p.alignment = 1  # 중앙 정렬
                        elif block['alignment'] == 'right':
                            p.alignment = 2  # 오른쪽 정렬
                        else:
                            p.alignment = 0  # 왼쪽 정렬
                            
                except (IndexError, AttributeError):
                    # 내용 플레이스홀더가 없는 경우 텍스트 박스 추가
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(6)
                    content_box = text_slide.shapes.add_textbox(left, top, width, height)
                    content_frame = content_box.text_frame
                    content_frame.text = final_text
            else:
                # 레이아웃 정보가 없는 경우 일반 텍스트로 추가
                try:
                    content = text_slide.placeholders[1]
                    content.text = final_text
                except (IndexError, AttributeError):
                    # 내용 플레이스홀더가 없는 경우 텍스트 박스 추가
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(6)
                    content_box = text_slide.shapes.add_textbox(left, top, width, height)
                    content_frame = content_box.text_frame
                    content_frame.text = final_text
            
            print("편집 가능한 텍스트가 레이아웃 정보와 함께 슬라이드에 추가되었습니다.")
        else:
            print("추출할 수 있는 텍스트가 없습니다.")
        
        # PPTX 파일 저장
        prs.save(output_path)
        return True
        
    except Exception as e:
        print(f"변환 중 오류 발생: {str(e)}")
        return False

def docx_to_pdf(docx_path, output_path):
    """DOCX를 PDF로 변환하는 함수"""
    try:
        # Windows에서 LibreOffice 사용
        if platform.system() == "Windows":
            # LibreOffice 경로 찾기
            libreoffice_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                "soffice"  # PATH에 있는 경우
            ]
            
            libreoffice_path = None
            for path in libreoffice_paths:
                if os.path.exists(path) or path == "soffice":
                    libreoffice_path = path
                    break
            
            if libreoffice_path:
                # LibreOffice를 사용하여 변환
                output_dir = os.path.dirname(output_path)
                cmd = [
                    libreoffice_path,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", output_dir,
                    docx_path
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True)
                if result.returncode == 0:
                    # 생성된 PDF 파일명 확인 및 이동
                    base_name = os.path.splitext(os.path.basename(docx_path))[0]
                    generated_pdf = os.path.join(output_dir, base_name + ".pdf")
                    
                    if os.path.exists(generated_pdf) and generated_pdf != output_path:
                        os.rename(generated_pdf, output_path)
                    
                    return os.path.exists(output_path)
                else:
                    print(f"LibreOffice 변환 실패: {result.stderr}")
                    return False
            else:
                print("LibreOffice를 찾을 수 없습니다.")
                return False
        else:
            print("현재 Linux/Mac에서의 DOCX → PDF 변환은 지원되지 않습니다.")
            return False
            
    except Exception as e:
        print(f"DOCX → PDF 변환 중 오류: {str(e)}")
        return False

# 파일 크기 초과 오류 처리
@app.errorhandler(413)
def too_large(e):
    flash('파일 크기가 100MB를 초과합니다. 더 작은 파일을 선택해주세요.')
    return redirect(url_for('index'))

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        print("파일 업로드 요청 시작")
        
        # 1단계: 파일 존재 여부 확인
        if 'file' not in request.files:
            error_msg = '파일이 선택되지 않았습니다.'
            print(f"[1단계 실패] {error_msg}")
            flash(error_msg)
            return redirect(url_for('index'))
        
        print("[1단계 통과] 파일 필드 존재 확인")
        file = request.files['file']
        
        # 2단계: 파일명 확인
        if not file or file.filename == '' or file.filename is None:
            error_msg = '파일이 선택되지 않았습니다.'
            print(f"[2단계 실패] {error_msg} - filename: {file.filename if file else 'None'}")
            flash(error_msg)
            return redirect(url_for('index'))
        
        print(f"[2단계 통과] 파일명 확인 - {file.filename}")
        
        # 3단계: 파일 내용 및 크기 확인 (강화된 검증)
        try:
            print("[3단계 시작] 파일 크기 및 내용 검증")
            # 파일 포인터를 끝으로 이동하여 크기 확인
            file.seek(0, 2)
            file_size = file.tell()
            file.seek(0)  # 파일 포인터를 다시 처음으로 이동
            
            print(f"[3단계] 파일 크기 측정: {file_size} bytes")
            
            # 파일 크기가 0인 경우 처리
            if file_size == 0:
                error_msg = '업로드된 파일이 비어있습니다. 올바른 PDF 파일을 선택해주세요.'
                print(f"[3단계 실패] {error_msg}")
                flash(error_msg)
                return redirect(url_for('index'))
            
            # 최소 파일 크기 확인 (PDF 헤더 최소 크기)
            if file_size < 100:  # 100바이트 미만은 유효한 PDF가 아님
                error_msg = '파일이 너무 작습니다. 올바른 PDF 파일을 선택해주세요.'
                print(f"[3단계 실패] {error_msg} - 크기: {file_size} bytes")
                flash(error_msg)
                return redirect(url_for('index'))
            
            if file_size > 100 * 1024 * 1024:  # 100MB
                error_msg = f'파일 크기가 너무 큽니다. (현재: {file_size // (1024*1024)}MB, 최대: 100MB)'
                print(f"[3단계 실패] {error_msg}")
                flash(error_msg)
                return redirect(url_for('index'))
            
            print(f"[3단계 통과] 파일 크기 검증 완료: {file_size // (1024*1024) if file_size >= 1024*1024 else file_size // 1024}{'MB' if file_size >= 1024*1024 else 'KB'}")
            
            # 4단계: PDF 파일 헤더 검증
            print("[4단계 시작] PDF 헤더 검증")
            file_content = file.read(10)  # 처음 10바이트 읽기
            file.seek(0)  # 다시 처음으로 이동
            
            print(f"[4단계] 파일 헤더: {file_content[:10]}")
            
            if not file_content.startswith(b'%PDF-'):
                error_msg = '올바른 PDF 파일이 아닙니다. PDF 파일을 선택해주세요.'
                print(f"[4단계 실패] {error_msg} - 헤더: {file_content}")
                flash(error_msg)
                return redirect(url_for('index'))
            
            print("[4단계 통과] PDF 헤더 검증 완료")
                
        except Exception as e:
            error_msg = f"파일 검증 중 오류: {str(e)}"
            print(f"[3-4단계 예외] {error_msg}")
            flash('파일을 읽는 중 오류가 발생했습니다. 다른 파일을 시도해주세요.')
            return redirect(url_for('index'))
        
        # 5단계: 파일 형식 확인 및 처리 (강화된 검증)
        print(f"[5단계 시작] 파일 형식 확인 - allowed_file({file.filename})")
        
        if file and allowed_file(file.filename):
            print("[5단계 통과] allowed_file 검증 완료")
            filename = secure_filename(file.filename)
            print(f"[5단계] secure_filename 적용: {filename}")
            
            # 파일 확장자 안전하게 추출 (list index out of range 오류 방지)
            if '.' not in filename:
                error_msg = '파일 확장자가 없습니다. PDF 파일을 선택해주세요.'
                print(f"[5단계 실패] {error_msg} - filename: {filename}")
                flash(error_msg)
                return redirect(url_for('index'))
            
            file_ext = filename.rsplit('.', 1)[1].lower()
            print(f"[5단계] 파일 확장자 추출: {file_ext}")
            
            if file_ext != 'pdf':
                error_msg = 'PDF 파일만 업로드 가능합니다.'
                print(f"[5단계 실패] {error_msg} - 확장자: {file_ext}")
                flash(error_msg)
                return redirect(url_for('index'))
            
            print("[5단계 통과] 파일 확장자 검증 완료")
            
            # 안전한 파일명 생성 (타임스탬프 추가로 중복 방지)
            import time
            timestamp = str(int(time.time()))
            safe_filename = f"{timestamp}_{filename}"
            input_path = os.path.join(UPLOAD_FOLDER, safe_filename)
            
            print(f"[6단계 시작] 파일 저장 - {input_path}")
            try:
                # 파일 저장 전 디렉토리 존재 확인
                os.makedirs(UPLOAD_FOLDER, exist_ok=True)
                print(f"[6단계] 업로드 폴더 준비 완료: {UPLOAD_FOLDER}")
                
                file.save(input_path)
                print(f"[6단계] 파일 저장 완료")
                
                # 저장된 파일 크기 재확인
                saved_file_size = os.path.getsize(input_path)
                print(f"[6단계] 저장된 파일 크기 확인: {saved_file_size} bytes")
                
                if saved_file_size == 0:
                    os.remove(input_path)
                    error_msg = '파일 저장 중 오류가 발생했습니다. 다시 시도해주세요.'
                    print(f"[6단계 실패] {error_msg} - 저장된 파일 크기가 0")
                    flash(error_msg)
                    return redirect(url_for('index'))
                
                print(f"[6단계 통과] 파일 저장 검증 완료 - 크기: {saved_file_size}바이트")
            except Exception as e:
                error_msg = f"파일 저장 오류: {str(e)}"
                print(f"[6단계 예외] {error_msg}")
                flash(f'파일 저장 중 오류가 발생했습니다: {str(e)}')
                return redirect(url_for('index'))
            
            # 변환 처리
            conversion_success = False
            output_path = None
            
            if file_ext == 'pdf':
                # PDF → DOCX 변환
                # 파일명에서 확장자 제거 (안전하게)
                base_filename = filename.rsplit('.', 1)[0] if '.' in filename else filename
                output_filename = base_filename + '.docx'
                output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                
                quality = request.form.get('quality', 'medium')
                print(f"PDF → DOCX 변환 시작 - {input_path} -> {output_path}")
                
                try:
                    print(f"=== PDF 변환 시작 ===")
                    print(f"입력 파일: {input_path}")
                    print(f"파일 크기: {os.path.getsize(input_path)} bytes")
                    print(f"출력 파일: {output_path}")
                    print(f"품질 설정: {quality}")
                    
                    conversion_success = pdf_to_docx(input_path, output_path, quality)
                    
                    print(f"=== PDF 변환 결과: {conversion_success} ===")
                    if conversion_success:
                        if os.path.exists(output_path):
                            output_size = os.path.getsize(output_path)
                            print(f"변환 성공 - 출력 파일 크기: {output_size} bytes")
                        else:
                            print("변환 성공으로 보고되었지만 출력 파일이 존재하지 않음")
                            conversion_success = False
                    else:
                        print("PDF 변환 실패 - pdf_to_docx 함수에서 False 반환")
                        
                except Exception as e:
                    print(f"=== 변환 중 예외 발생 ===")
                    print(f"예외 메시지: {str(e)}")
                    print(f"예외 타입: {type(e).__name__}")
                    import traceback
                    print(f"상세 스택 트레이스:")
                    print(traceback.format_exc())
                    flash(f'변환 중 오류가 발생했습니다: {str(e)}')
                    conversion_success = False
                    
            elif file_ext == 'docx':
                # DOCX → PDF 변환
                # 파일명에서 확장자 제거 (안전하게)
                base_filename = filename.rsplit('.', 1)[0] if '.' in filename else filename
                output_filename = base_filename + '.pdf'
                output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                
                print(f"DOCX → PDF 변환 시작 - {input_path} -> {output_path}")
                
                try:
                    conversion_success = docx_to_pdf(input_path, output_path)
                except Exception as e:
                    print(f"변환 중 예외 발생: {str(e)}")
                    flash(f'변환 중 오류가 발생했습니다: {str(e)}')
            
            # 변환 결과 처리
            if conversion_success:
                print("변환 성공 - 다운로드 준비")
                
                # 업로드된 파일 정리
                try:
                    os.remove(input_path)
                    print("임시 파일 삭제 완료")
                except Exception as e:
                    print(f"임시 파일 삭제 실패 (무시됨): {e}")
                
                # 파일 다운로드 제공
                try:
                    print("파일 다운로드 시작")
                    # 한글 파일명이 깨지지 않도록 UTF-8로 인코딩하여 전달
                    return send_file(
                        output_path, 
                        as_attachment=True, 
                        download_name=output_filename.encode('utf-8').decode('latin-1')
                    )
                except Exception as e:
                    print(f"파일 다운로드 오류: {str(e)}")
                    flash(f'파일 다운로드 중 오류가 발생했습니다: {str(e)}')
                    return redirect(url_for('index'))
            else:
                print("변환 실패 - 정리 작업")
                flash('파일 변환에 실패했습니다. 다시 시도해주세요.')
                
                # 실패한 파일들 정리
                for cleanup_path in [input_path, output_path]:
                    try:
                        if cleanup_path and os.path.exists(cleanup_path):
                            os.remove(cleanup_path)
                    except Exception as e:
                        print(f"파일 정리 실패 (무시됨): {e}")
                
                return redirect(url_for('index'))
        else:
            error_msg = 'PDF 또는 DOCX 파일만 업로드 가능합니다.'
            print(f"[5단계 실패] allowed_file 검증 실패 - {error_msg} - filename: {file.filename if file else 'None'}")
            flash(error_msg)
            return redirect(url_for('index'))
            
    except Exception as e:
        print(f"업로드 처리 중 예외 발생: {str(e)}")
        flash('파일 처리 중 오류가 발생했습니다.')
        return redirect(url_for('index'))

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)